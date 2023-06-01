import requests
import urllib.request
from bs4 import BeautifulSoup
import os
from pptx import Presentation
import os
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# insert link to the file server where the files are accessable
url = ""

response = requests.get(url)
html_content = response.text

# BeautifulSoup for analyzing the HTML-structure and detect the files
soup = BeautifulSoup(html_content, "html.parser")
links = []

# saving all links to the files
for link in soup.find_all("a"):
    href = link.get("href")
    if href.endswith(".txt") or href.endswith(".jpg"):
        links.append(href)

# urllib for downloading the files
for link in links:
    file_url = os.path.join(url, link)
    file_name = os.path.basename(link)
    urllib.request.urlretrieve(file_url, file_name)

prs = Presentation('template.pptx')

# add all majors - 0 is the counter
majors = {
    "IF": ["Bachelor Informatik", 0], 
    "IB": ["Bachelor Wirtschaftsinformatik", 0], 
    "DC": ["Bachelor Data Science & Scientific Computing", 0],
    "IT": ["Master Informatik", 0], 
    "IS": ["Master IT-Security", 0], 
    "IN": ["Master Wirtschaftinformatik", 0]
    }

name_list = []

for major in majors: 
    slide_layout = prs.slides[5]
    new_slide = prs.slides.add_slide(slide_layout.slide_layout)
    for shape in new_slide.placeholders:
        shape.text_frame.text = majors[major][0]
    
    for file_name in os.listdir("./"):
        if file_name.endswith(".txt"):
            txt_name = file_name.removesuffix('.txt')    
            with open(os.path.join(file_name), 'r') as file:
                contents = file.readlines()[0]

                # removing tabs and "categories"
                contents = contents.replace("\t", "")
                contents = contents.replace("\n", "")
                contents = contents.replace("Nachname: ", "|")
                contents = contents.replace("Vorname: ", "|")
                contents = contents.replace("Highlight:", "|")
                contents = contents.replace("Hobbies:", "|")
                contents = contents.replace("Studiengang: ", "|")
                contents = contents.split("|")[1:]

            if major not in contents:
                continue
            else:
                majors[major][1] += 1
            
            picture = os.path.isfile(txt_name + ".jpg")

            last_name = contents[0]
            first_name = contents[1]

            #name_list.append(last_name + " " + first_name)

            if picture:
                # only picture and name is avalable
                if len(contents[2].strip()) == 0 and len(contents[3].strip()) == 0:
                    slide_layout1 = prs.slides[4]
                    new_slide = prs.slides.add_slide(slide_layout1.slide_layout)#, index=9) 
                    for shape in new_slide.placeholders:
                        if shape.name == "Text Placeholder 2":
                            shape.text_frame.text = first_name
                        elif shape.name == "Text Placeholder 3":
                            shape.text_frame.text = last_name
                        elif shape.name == "Picture Placeholder 1":
                            picture = shape.insert_picture(txt_name + ".jpg")

                # picture, name, and hobby/highlight is available       
                else:
                    highlights = contents[2]
                    hobbies = contents[3]
                    studies = contents[4]
                    slide_layout1 = prs.slides[1]
                    new_slide = prs.slides.add_slide(slide_layout1.slide_layout)#, index=9) 
                    for shape in new_slide.placeholders:
                        if shape.name == "Text Placeholder 4":
                            shape.text_frame.text = first_name
                        elif shape.name == "Text Placeholder 5":
                            shape.text_frame.text = last_name
                        elif shape.name == "Text Placeholder 2":
                            shape.text_frame.text = highlights
                        elif shape.name == "Text Placeholder 3":
                            shape.text_frame.text = hobbies
                        elif shape.name == "Picture Placeholder 1":
                            picture = shape.insert_picture(txt_name + ".jpg")
            else:
                # only name is available
                if len(contents[2].strip()) == 0 and len(contents[3].strip()) == 0:
                    slide_layout1 = prs.slides[2]
                    new_slide = prs.slides.add_slide(slide_layout1.slide_layout)#, index=9) 
                    for shape in new_slide.placeholders:
                        if shape.name == "Text Placeholder 1":
                            shape.text_frame.text = first_name
                        elif shape.name == "Text Placeholder 2":
                            shape.text_frame.text = last_name
                # name and hobby/highlight is available
                else:
                    slide_layout1 = prs.slides[3]
                    new_slide = prs.slides.add_slide(slide_layout1.slide_layout)#, index=9) 
                    for shape in new_slide.placeholders:
                        if shape.name == "Text Placeholder 3":
                            shape.text_frame.text = first_name
                        elif shape.name == "Text Placeholder 4":
                            shape.text_frame.text = last_name
                        elif shape.name == "Text Placeholder 1":
                            shape.text_frame.text = highlights
                        elif shape.name == "Text Placeholder 2":
                            shape.text_frame.text = hobbies

prs.save('output.pptx')

for info in majors:
    print(majors[info])